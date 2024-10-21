from typing import List
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import os
import glob
from tqdm import tqdm
from multiprocessing import Pool
import logging

# 定义源目录和块参数
source_directory = "/root/ChatChat/books"
chunk_size = 1000
chunk_overlap = 200

logging.basicConfig(level=logging.DEBUG)


# 从PPT文件加载文本的函数
def load_ppt(file_path: str) -> List[str]:
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return text


# 从DOCX文件加载文本的函数
def load_docx(file_path: str) -> List[str]:
    docx = DocxDocument(file_path)
    text = [paragraph.text for paragraph in docx.paragraphs if paragraph.text.strip()]
    return text


# 从CSV文件加载文本的函数
def load_csv(file_path: str) -> List[str]:
    df = pd.read_csv(file_path)
    return df.to_string(index=False).split('\n')


# 从XLSX文件加载文本的函数
def load_xlsx(file_path: str) -> List[str]:
    df = pd.read_excel(file_path, engine='openpyxl')
    return df.to_string(index=False).split('\n')


# 从TXT文件加载文本的函数
def load_txt(file_path: str) -> List[str]:
    with open(file_path, 'r', encoding='utf8') as f:
        text = f.readlines()
    return text


# 加载单个文档的函数
def load_single_document(file_path: str) -> List[str]:
    # logging.debug(f"加载文档: {file_path}")
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == '.ppt' or ext == '.pptx':
        return load_ppt(file_path)
    elif ext == '.docx':
        return load_docx(file_path)
    elif ext == '.csv':
        return load_csv(file_path)
    elif ext == '.xlsx' or ext == '.xls':
        return load_xlsx(file_path)
    elif ext == '.txt':
        return load_txt(file_path)
    else:
        raise ValueError(f"不支持的文件扩展名 '{ext}'")


# 从目录中加载所有文档的函数
def load_documents(source_dir: str, ignored_files: List[str] = None) -> List[List[str]]:
    ignored_files = [] if ignored_files is None else ignored_files
    all_files = []
    for ext in ['.ppt', '.pptx', '.docx', '.csv', '.xlsx', '.xls', '.txt']:
        all_files.extend(glob.glob(os.path.join(source_dir, f"**/*{ext}"), recursive=True))
    filtered_files = [file_path for file_path in all_files if file_path not in ignored_files]

    documents = []
    with Pool(processes=os.cpu_count()) as pool:
        with tqdm(total=len(filtered_files), desc='加载新文档', ncols=80) as pbar:
            for docs in pool.imap_unordered(load_single_document, filtered_files):
                documents.append(docs)
                pbar.update()

    return documents


# 将文本拆分为块的函数
def split_into_chunks(text: List[str], chunk_size: int, chunk_overlap: int) -> List[str]:
    chunks = []
    current_chunk = []
    current_size = 0

    for line in text:
        line_length = len(line)
        if current_size + line_length < chunk_size:
            current_chunk.append(line)
            current_size += line_length
        else:
            if current_chunk:
                chunks.append("\n".join(current_chunk))
            current_chunk = current_chunk[-chunk_overlap:]  # 保留重叠部分
            current_size = sum(len(l) for l in current_chunk)
            current_chunk.append(line)
            current_size += line_length

    if current_chunk:
        chunks.append("\n".join(current_chunk))

    return chunks


# 处理文档的函数
def process_documents(ignored_files: List[str] = None) -> List[str]:
    ignored_files = [] if ignored_files is None else ignored_files
    print(f"从 {source_directory} 加载文档")
    documents = load_documents(source_directory, ignored_files)

    if not documents:
        print("没有新文档可加载")
        import sys
        sys.exit(0)

    print(f"从 {source_directory} 加载了 {len(documents)} 个新文档")

    all_chunks = []
    for document_text in documents:
        document_chunks = split_into_chunks(document_text, chunk_size, chunk_overlap)
        all_chunks.extend(document_chunks)

    print(f"拆分成 {len(all_chunks)} 个文本块（最大 {chunk_size} 个字符）")
    return all_chunks


if __name__ == "__main__":
    document_chunks = process_documents()
    print("生成的文本块:")
    for i, chunk in enumerate(document_chunks[:5]):
        print(f"文本块 {i + 1}:\n{chunk}\n")
